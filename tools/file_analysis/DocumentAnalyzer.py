from typing import Dict, Any, List, Optional
from ..base import BaseTool
from pathlib import Path
import logging
import os
import io

# 引入必要的分析库
# 注意：你需要确保环境安装了: pymupdf, python-docx, python-pptx, pandas, openpyxl
import fitz  # PyMuPDF
import pandas as pd
try:
    import docx
except ImportError:
    docx = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from tools.base import ToolResult, ToolFailure

from utils.logger import logger

class DocumentAnalyzerTool(BaseTool):
    """
    全能文件分析工具 (支持 PDF, Word, Excel, CSV, PPT)
    基于 'Lunisia' 文件处理能力增强，支持深度内容提取与统计分析。
    """

    name: str = "analyze_document"
    description: str = (
        "🔍 专业文件内容分析工具。用于读取和深度分析各类文档，支持 PDF、Word、Excel、CSV、PPT 等格式。"
        "返回完整文本内容、元数据、结构化数据统计。任何需要查看或分析文件内容的场景都应使用此工具。"
    )
    parameters: dict = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "文件的本地绝对路径"
            },
            "preview_length": {
                "type": "integer",
                "description": "文本预览截取的字符长度，默认为 10000",
                "default": 10000
            }
        },
        "required": ["file_path"]
    }

    async def execute(self, file_path: str, preview_length: int = 10000) -> ToolResult:
        path = Path(file_path)
        if not path.exists():
            return ToolFailure(error=f"文件不存在: {file_path}")

        suffix = path.suffix.lower()
        file_name = path.name

        try:
            # 根据后缀分发处理逻辑
            if suffix == ".pdf":
                result = self._analyze_pdf(str(path))
            elif suffix in [".docx", ".doc"]:
                result = self._analyze_docx(str(path))
            elif suffix in [".pptx", ".ppt"]:
                result = self._analyze_pptx(str(path))
            elif suffix in [".csv", ".xlsx", ".xls"]:
                result = self._analyze_table(str(path), suffix)
            else:
                return ToolFailure(error=f"不支持的文件格式: {suffix}\n请检查文件格式并重新上传。\n支持的格式包括: PDF, Word (.docx), PowerPoint (.pptx), Excel (.xlsx, .xls), CSV.")

            # 构造返回给 LLM 的最终结果
            # System 部分放置元数据和结构化分析，Output 放置具体文本内容
            metadata_str = "\n".join([f"- {k}: {v}" for k, v in result['metadata'].items()])
            analysis_str = result.get('analysis', '无额外分析')
            
            full_content = result['content']
            # 如果内容过长，进行截断处理，但在 system 中提示
            display_content = full_content[:preview_length]
            if len(full_content) > preview_length:
                display_content += f"\n\n[...剩余内容已截断，总长度 {len(full_content)} 字符...]"

            system_msg = (
                f"文件分析报告: {file_name}\n"
                f"【元数据】\n{metadata_str}\n"
                f"【智能分析】\n{analysis_str}"
            )

            return ToolResult(output=display_content, system=system_msg)

        except Exception as e:
            logger.error(f"文件分析失败: {e}", exc_info=True)
            return ToolFailure(error=f"分析过程中发生错误: {str(e)}")

    def _analyze_pdf(self, file_path: str) -> Dict[str, Any]:
        """PDF深度分析 (参考 file_analysis_tool.PDFAnalyzer)"""
        doc = fitz.open(file_path)
        content_parts = []
        
        metadata = {
            "type": "PDF Document",
            "page_count": doc.page_count,
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "creation_date": doc.metadata.get("creationDate", ""),
        }

        # 提取文本
        for page_num, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                content_parts.append(f"--- 第 {page_num + 1} 页 ---\n{text}")
        
        full_text = "\n".join(content_parts)
        
        # 简单分析：检测是否有表格或代码特征
        analysis = []
        if "Table" in full_text or "表" in full_text:
            analysis.append("文档可能包含表格数据。")
        if len(full_text) > 0:
            analysis.append(f"有效文本长度: {len(full_text)} 字符。")
        
        doc.close()
        return {
            "metadata": metadata,
            "content": full_text,
            "analysis": "\n".join(analysis)
        }

    def _analyze_docx(self, file_path: str) -> Dict[str, Any]:
        """Word深度分析 (参考 file_analysis_tool.DocxAnalyzer)"""
        if docx is None:
            raise ImportError("请安装 python-docx 以分析 Word 文档")
            
        doc = docx.Document(file_path)
        core_props = doc.core_properties
        
        metadata = {
            "type": "Word Document",
            "title": core_props.title or "Unknown",
            "author": core_props.author or "Unknown",
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
            "modified": str(core_props.modified) if core_props.modified else ""
        }

        content_parts = []
        
        # 提取段落
        for para in doc.paragraphs:
            if para.text.strip():
                content_parts.append(para.text)
        
        # 提取表格内容 (这是参考文件中非常有用的功能)
        if doc.tables:
            content_parts.append("\n--- 文档内表格数据 ---")
            for i, table in enumerate(doc.tables):
                content_parts.append(f"[表格 {i+1}]")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    content_parts.append(row_text)

        full_text = "\n\n".join(content_parts)
        
        # 结构分析
        analysis = []
        if metadata['table_count'] > 0:
            analysis.append(f"包含 {metadata['table_count']} 个表格，涉及结构化数据。")
        
        # 简单的标题检测
        headings = sum(1 for p in doc.paragraphs if p.style.name.startswith('Heading'))
        if headings > 0:
            analysis.append(f"检测到 {headings} 个标题层级，文档结构清晰。")

        return {
            "metadata": metadata,
            "content": full_text,
            "analysis": "\n".join(analysis)
        }

    def _analyze_pptx(self, file_path: str) -> Dict[str, Any]:
        """PPTX分析"""
        if Presentation is None:
            raise ImportError("请安装 python-pptx 以分析 PPT 文档")

        prs = Presentation(file_path)
        metadata = {
            "type": "PowerPoint Presentation",
            "slide_count": len(prs.slides)
        }
        
        content_parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # 提取标题
            if slide.shapes.title:
                slide_text.append(f"Title: {slide.shapes.title.text}")
            
            # 提取文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    if shape.text.strip():
                        slide_text.append(shape.text)
            
            content_parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

        return {
            "metadata": metadata,
            "content": "\n\n".join(content_parts),
            "analysis": f"共包含 {len(prs.slides)} 张幻灯片。"
        }

    def _analyze_table(self, file_path: str, suffix: str) -> Dict[str, Any]:
        """表格数据分析 (参考 file_analysis_tool.TableAnalyzer)"""
        if suffix == '.csv':
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
            
        rows, cols = df.shape
        
        metadata = {
            "type": "Table Data (CSV/Excel)",
            "rows": rows,
            "columns": cols,
            "column_names": df.columns.tolist(),
            "memory_usage": f"{df.memory_usage(deep=True).sum() / 1024:.2f} KB"
        }

        # 生成数据预览（Markdown 格式）
        # 限制行数以避免 Token 爆炸
        preview_rows = min(rows, 20) 
        content = f"前 {preview_rows} 行数据预览:\n"
        content += df.head(preview_rows).to_markdown(index=False)

        # 深度统计分析 (参考 Lunisia 的 TableAnalyzer)
        analysis_parts = []
        
        # 1. 缺失值统计
        missing = df.isnull().sum()
        if missing.sum() > 0:
            missing_cols = missing[missing > 0]
            missing_info = ", ".join([f"{col}({val})" for col, val in missing_cols.items()])
            analysis_parts.append(f"⚠️ 缺失值检测: {missing_info}")
        else:
            analysis_parts.append("✅ 数据完整（无缺失值）。")

        # 2. 数据类型推断
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        if numeric_cols:
            analysis_parts.append(f"📈 数值列 ({len(numeric_cols)}个): {', '.join(numeric_cols[:5])}...")
            # 简单的数值统计
            stats = df[numeric_cols].describe().to_markdown()
            content += f"\n\n数值列统计描述:\n{stats}"

        # 3. 时间列检测
        time_candidates = [col for col in df.columns if 'date' in str(col).lower() or 'time' in str(col).lower()]
        if time_candidates:
            analysis_parts.append(f"⏰ 可能的时间列: {', '.join(time_candidates)}")

        return {
            "metadata": metadata,
            "content": content,
            "analysis": "\n".join(analysis_parts)
        }